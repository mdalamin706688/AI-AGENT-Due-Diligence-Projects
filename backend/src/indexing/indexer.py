from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.embeddings import Embeddings
import numpy as np
import hashlib
from ..models import Document
from ..storage.memory import storage
import uuid
from pypdf import PdfReader
import os
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

class MockEmbeddings(Embeddings):
    """Mock embeddings for demo purposes"""
    
    def embed_query(self, text: str) -> list:
        # Create a deterministic hash-based embedding
        hash_obj = hashlib.md5(text.encode())
        hash_bytes = hash_obj.digest()
        # Convert to float array and normalize
        embedding = np.array([b / 255.0 for b in hash_bytes])
        # Pad to 1536 dimensions (OpenAI ada-002 dimension)
        embedding = np.pad(embedding, (0, 1536 - len(embedding)), 'constant')
        return embedding.tolist()
    
    def embed_documents(self, texts: list) -> list:
        return [self.embed_query(text) for text in texts]

class DocumentIndexer:
    def __init__(self):
        # Use mock embeddings for demo
        self.embeddings = MockEmbeddings()
        self.vectorstore = None
        self.documents_indexed = set()

    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from various file formats"""
        if file_path.endswith('.pdf'):
            return self.extract_text_from_pdf(file_path)
        elif file_path.endswith('.txt'):
            return self.extract_text_from_txt(file_path)
        elif file_path.endswith('.docx'):
            return self.extract_text_from_docx(file_path)
        elif file_path.endswith('.xlsx') or file_path.endswith('.xls'):
            return self.extract_text_from_excel(file_path)
        elif file_path.endswith('.pptx') or file_path.endswith('.ppt'):
            return self.extract_text_from_powerpoint(file_path)
        else:
            print(f"Unsupported file type: {file_path}")
            return ""

    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading TXT file {file_path}: {e}")
            return ""

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            print(f"Error reading PDF file {file_path}: {e}")
            return ""

    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            return text
        except Exception as e:
            print(f"Error reading DOCX file {file_path}: {e}")
            return ""

    def extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file (XLSX/XLS)"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text
        except Exception as e:
            print(f"Error reading Excel file {file_path}: {e}")
            return ""

    def extract_text_from_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint file (PPTX/PPT)"""
        try:
            presentation = Presentation(file_path)
            text = ""
            for slide_number, slide in enumerate(presentation.slides, 1):
                text += f"Slide {slide_number}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                text += "\n"
            return text
        except Exception as e:
            print(f"Error reading PowerPoint file {file_path}: {e}")
            return ""

    def index_document(self, doc: Document):
        """Index a document into the vector store"""
        if doc.id in self.documents_indexed:
            return  # Already indexed
        
        # Extract text if not already done
        if not doc.content.strip():
            if os.path.exists(doc.filename):
                doc.content = self.extract_text_from_file(doc.filename)
            else:
                print(f"Document file not found: {doc.filename}")
                return
        
        if not doc.content.strip():
            print(f"No content extracted from {doc.filename}")
            return
        
        # Split text into chunks with better strategy for headers
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,  # Larger chunks to preserve more context
            chunk_overlap=300,  # More overlap to maintain continuity
            separators=["\n\n\n", "\n\n", "\n", ". ", " ", ""]  # Prioritize paragraph breaks
        )
        chunks = text_splitter.split_text(doc.content)
        
        # Create chunks with metadata
        doc.chunks = []
        for i, chunk in enumerate(chunks):
            if chunk.strip():  # Only add non-empty chunks
                chunk_id = str(uuid.uuid4())
                doc.chunks.append({
                    "id": chunk_id,
                    "text": chunk,
                    "metadata": {
                        "doc_id": doc.id,
                        "chunk_index": i,
                        "filename": doc.filename
                    }
                })
        
        if not doc.chunks:
            print(f"No chunks created for {doc.filename}")
            return
        
        # Add to vector store
        texts = [c["text"] for c in doc.chunks]
        metadatas = [c["metadata"] for c in doc.chunks]
        
        if self.vectorstore is None:
            self.vectorstore = FAISS.from_texts(texts, self.embeddings, metadatas=metadatas)
        else:
            self.vectorstore.add_texts(texts, metadatas=metadatas)
        
        self.documents_indexed.add(doc.id)
        storage.save_document(doc)
        print(f"Indexed document {doc.filename} with {len(doc.chunks)} chunks")

    def search(self, query: str, k=5, document_ids=None):
        """Search the vector store for relevant chunks with improved keyword mapping"""
        query_lower = query.lower()
        
        # For questions that need specific factual information, prioritize keyword search
        question_keywords = ['what is', 'what are', 'who is', 'how much', 'how many', 'where', 'when']
        is_specific_question = any(qk in query_lower for qk in question_keywords)
        
        if is_specific_question:
            # Use keyword search first for specific questions
            return self.enhanced_keyword_search(query, k=k, document_ids=document_ids)
        
        # First try semantic search
        if self.vectorstore:
            semantic_results = self.vectorstore.similarity_search(query, k=k)
            # Filter by document_ids if provided
            if document_ids:
                semantic_results = [r for r in semantic_results if r.metadata.get('doc_id') in document_ids]
            
            # Check if semantic results are actually relevant (contain query keywords)
            query_words = set(query.lower().split())
            relevant_semantic = []
            for result in semantic_results:
                result_text = result.page_content.lower()
                if any(word in result_text for word in query_words):
                    relevant_semantic.append(result)
            
            if relevant_semantic:
                return relevant_semantic[:k]
        
        # Fallback to enhanced keyword search
        return self.enhanced_keyword_search(query, k=k, document_ids=document_ids)
    
    def enhanced_keyword_search(self, query: str, k=5, document_ids=None):
        """Enhanced keyword search with question-to-content mapping"""
        from ..storage.memory import storage
        
        query_lower = query.lower()
        
        # Map common question patterns to likely content keywords
        keyword_mappings = {
            'company name': ['company', 'corporation', 'inc', 'ltd', 'group', 'incorporated'],
            'revenue': ['revenue', 'sales', 'income', 'turnover', 'earnings'],
            'financial statements': ['financial', 'statement', 'balance sheet', 'income statement', 'cash flow'],
            'lawsuits': ['lawsuit', 'litigation', 'legal', 'court', 'claim'],
            'business model': ['business model', 'operations', 'strategy', 'services'],
            'key financial metrics': ['metrics', 'kpi', 'performance', 'ratios', 'valuation']
        }
        
        # Find matching keywords for this query
        search_keywords = set()
        for question_pattern, content_keywords in keyword_mappings.items():
            if question_pattern in query_lower:
                search_keywords.update(content_keywords)
        
        # If no specific mapping found, use original query words
        if not search_keywords:
            search_keywords = set(query_lower.split())
        
        # Also include original query terms
        search_keywords.update(query_lower.split())
        
        results = []
        
        # Search through indexed documents, filtered by document_ids if provided
        docs_to_search = self.documents_indexed
        if document_ids:
            docs_to_search = docs_to_search.intersection(set(document_ids))
        
        for doc_id in docs_to_search:
            doc = storage.get_document(doc_id)
            if not doc or not doc.chunks:
                continue
                
            for chunk_data in doc.chunks:
                chunk_text = chunk_data["text"].lower()
                # Count matching keywords (weighted by importance)
                score = 0
                matched_terms = []
                
                for keyword in search_keywords:
                    if keyword in chunk_text:
                        # Give higher weight to exact matches and company-related terms
                        weight = 2 if keyword in ['company', 'corporation', 'inc', 'group'] else 1
                        score += weight
                        matched_terms.append(keyword)
                
                if score > 0:
                    # Create a mock result object similar to FAISS results
                    class MockResult:
                        def __init__(self, page_content, metadata, score):
                            self.page_content = page_content
                            self.metadata = metadata
                            self.score = score
                    
                    results.append(MockResult(
                        page_content=chunk_data["text"],
                        metadata=chunk_data["metadata"],
                        score=score
                    ))
        
        # Sort by score (relevance) and return top k
        results.sort(key=lambda x: x.score, reverse=True)
        return results[:k]

    def keyword_search(self, query: str, k=5, document_ids=None):
        """Simple keyword-based search as fallback"""
        from ..storage.memory import storage
        
        query_lower = query.lower()
        query_words = set(query_lower.split())
        results = []
        
        # Search through indexed documents, filtered by document_ids if provided
        docs_to_search = self.documents_indexed
        if document_ids:
            docs_to_search = docs_to_search.intersection(set(document_ids))
        
        for doc_id in docs_to_search:
            doc = storage.get_document(doc_id)
            if not doc or not doc.chunks:
                continue
                
            for chunk_data in doc.chunks:
                chunk_text = chunk_data["text"].lower()
                # Count matching words
                matching_words = sum(1 for word in query_words if word in chunk_text)
                if matching_words > 0:
                    # Create a mock result object similar to FAISS results
                    class MockResult:
                        def __init__(self, page_content, metadata):
                            self.page_content = page_content
                            self.metadata = metadata
                    
                    results.append(MockResult(
                        page_content=chunk_data["text"],
                        metadata=chunk_data["metadata"]
                    ))
        
        # Sort by relevance (simple word count) and return top k
        results.sort(key=lambda x: len([w for w in query_words if w in x.page_content.lower()]), reverse=True)
        return results[:k]

indexer = DocumentIndexer()